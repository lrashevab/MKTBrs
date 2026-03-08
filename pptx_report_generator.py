# -*- coding: utf-8 -*-
"""
V2 PPTX 報告生成器 + Playwright PDF 生成器
- report_executive_summary.pptx  → 可直接匯入 Canva / Google 簡報（可編輯）
- report_battlecard.pptx         → 競品對決 Battlecard 投影片（可編輯）
- report_*.pdf                   → 快速列印用（Playwright 自動生成）
"""

import io
import json
import sys
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Optional

from output_utils import _safe_print


def _fix_win_utf8():
    if sys.platform != 'win32':
        return
    try:
        if hasattr(sys.stdout, 'buffer'):
            sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
    except Exception:
        pass
    try:
        if hasattr(sys.stderr, 'buffer'):
            sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8', errors='replace')
    except Exception:
        pass

# ── 顏色常數 ──────────────────────────────────────────
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

NAVY    = RGBColor(0x1a, 0x2f, 0x4b)
BLUE    = RGBColor(0x2d, 0x6a, 0x9f)
LIGHT_BLUE = RGBColor(0xe8, 0xf0, 0xfd)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY    = RGBColor(0x7f, 0x8c, 0x8d)
DARK    = RGBColor(0x2c, 0x3e, 0x50)
BG      = RGBColor(0xf0, 0xf4, 0xf8)
GREEN   = RGBColor(0x27, 0xae, 0x60)
RED     = RGBColor(0xe7, 0x4c, 0x3c)
ORANGE  = RGBColor(0xe6, 0x7e, 0x22)
YELLOW_BG = RGBColor(0xff, 0xf3, 0xcd)

# 16:9 widescreen
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ═══════════════════════════════════════════════════════
#  輔助函式
# ═══════════════════════════════════════════════════════

def _solid_bg(slide, color: RGBColor):
    """設定投影片純色背景。"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, left, top, width, height,
                 text: str, font_size: int = 14,
                 bold: bool = False,
                 color: RGBColor = DARK,
                 align=PP_ALIGN.LEFT,
                 wrap: bool = True) -> object:
    """新增文字框，回傳 text_frame。"""
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tf


def _add_rect(slide, left, top, width, height,
              fill_color: RGBColor,
              line_color: Optional[RGBColor] = None):
    """新增填色矩形。"""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape


def _header_slide(prs: Presentation, title: str,
                  subtitle: str, meta: str) -> None:
    """封面投影片（深藍漸層）。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, NAVY)

    # 左側色條
    _add_rect(slide, 0, 0, Inches(0.15), SLIDE_H, BLUE)

    # 主標題
    _add_textbox(slide,
                 Inches(0.5), Inches(2.2),
                 Inches(12), Inches(1.5),
                 title, font_size=36, bold=True,
                 color=WHITE, align=PP_ALIGN.LEFT)

    # 副標題
    _add_textbox(slide,
                 Inches(0.5), Inches(3.9),
                 Inches(11), Inches(0.8),
                 subtitle, font_size=16,
                 color=RGBColor(0xb8, 0xcf, 0xe8),
                 align=PP_ALIGN.LEFT)

    # 日期 meta
    _add_textbox(slide,
                 Inches(0.5), Inches(6.5),
                 Inches(10), Inches(0.5),
                 meta, font_size=11,
                 color=RGBColor(0x88, 0xaa, 0xcc),
                 align=PP_ALIGN.LEFT)


def _section_title_slide(prs: Presentation, number: str,
                         title: str, color: RGBColor = BLUE) -> None:
    """章節分隔投影片。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, color)
    _add_textbox(slide,
                 Inches(1), Inches(3),
                 Inches(11), Inches(1.2),
                 f'{number}  {title}',
                 font_size=32, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════
#  資料載入
# ═══════════════════════════════════════════════════════

def _load_data(campaign_dir: Path) -> Dict:
    d: Dict = {}

    def _read_json(path):
        try:
            return json.loads(path.read_text('utf-8'))
        except Exception:
            return None

    d['config'] = _read_json(campaign_dir / 'project_config.json') or {}

    raw = _read_json(campaign_dir / 'competitors.json')
    if isinstance(raw, dict) and 'competitors' in raw:
        d['competitors'] = raw['competitors']
    elif isinstance(raw, list):
        d['competitors'] = raw
    else:
        d['competitors'] = []

    kw = _read_json(campaign_dir / 'confirmed_keywords.json') or {}
    d['keywords'] = kw

    # Article stats from CSVs
    import csv
    def _read_csv(path):
        try:
            with open(path, 'r', encoding='utf-8-sig') as f:
                reader = csv.reader(f)
                headers = next(reader, [])
                rows = [r for r in reader if any(r)]
            return headers, rows
        except Exception:
            return [], []

    dcard_files = sorted(campaign_dir.glob('dcard_*.csv'))
    d['dcard_headers'], d['dcard_rows'] = (
        _read_csv(dcard_files[0]) if dcard_files else ([], [])
    )
    ptt_files = sorted(campaign_dir.glob('ptt_*.csv'))
    d['ptt_headers'], d['ptt_rows'] = (
        _read_csv(ptt_files[0]) if ptt_files else ([], [])
    )
    news_files = sorted(campaign_dir.glob('google_news_*.csv'))
    d['news_headers'], d['news_rows'] = (
        _read_csv(news_files[0]) if news_files else ([], [])
    )

    # AI analysis sections
    ana = campaign_dir / 'stage5_deep_analysis.md'
    d['deep_analysis_md'] = ana.read_text('utf-8') if ana.exists() else ''

    return d


def _top_articles(data: Dict, n: int = 5) -> List[Dict]:
    arts: List[Dict] = []
    for row in data.get('dcard_rows', []):
        h = data.get('dcard_headers', [])
        d = dict(zip(h, row)) if len(row) >= len(h) else {}
        try:
            eng = float(str(d.get('互動指數', 0)).replace(',', '') or 0)
        except Exception:
            eng = 0
        arts.append({'platform': 'Dcard',
                     'title': (d.get('標題', '') or '')[:50],
                     'engagement': eng,
                     'board': d.get('看板', '')})

    for row in data.get('ptt_rows', []):
        h = data.get('ptt_headers', [])
        d = dict(zip(h, row)) if len(row) >= len(h) else {}
        p = str(d.get('推文數', '0'))
        eng = 100 if p == '爆' else (-10 if p.startswith('X') else
                                     (int(p) if p.isdigit() else 0))
        arts.append({'platform': 'PTT',
                     'title': (d.get('標題', '') or '')[:50],
                     'engagement': eng,
                     'board': d.get('看板', '')})

    arts.sort(key=lambda x: x['engagement'], reverse=True)
    return arts[:n]


def _kw_count(data: Dict) -> int:
    kd = data.get('keywords', {})
    if kd.get('total_count'):
        return int(kd['total_count'])
    kw_obj = kd.get('keywords', kd)
    if isinstance(kw_obj, dict):
        return sum(len(v) for v in kw_obj.values() if isinstance(v, list))
    return 0


def _parse_md_sections(md: str) -> Dict[str, str]:
    sections: Dict[str, str] = {}
    current = None
    buf: List[str] = []
    for line in md.split('\n'):
        if line.startswith('## '):
            if current is not None:
                sections[current] = '\n'.join(buf).strip()
            current = line[3:].strip()
            buf = []
        elif current is not None:
            buf.append(line)
    if current is not None:
        sections[current] = '\n'.join(buf).strip()
    return sections


# ═══════════════════════════════════════════════════════
#  PPTX 生成：執行摘要版
# ═══════════════════════════════════════════════════════

def generate_executive_pptx(campaign_dir: Path, data: Dict) -> str:
    """生成執行摘要投影片（可匯入 Canva / Google 簡報）。"""
    cfg = data.get('config', {})
    brand   = cfg.get('brand', '品牌')
    service = cfg.get('service', '產品')
    industry = cfg.get('industry', '產業')
    purpose  = cfg.get('purpose', '競品分析')
    audience = cfg.get('audience', '目標受眾')

    competitors = data.get('competitors', [])
    dcard_cnt = len(data.get('dcard_rows', []))
    ptt_cnt   = len(data.get('ptt_rows', []))
    news_cnt  = len(data.get('news_rows', []))
    total_cnt = dcard_cnt + ptt_cnt + news_cnt
    kw_cnt    = _kw_count(data)
    top_arts  = _top_articles(data, 5)

    # AI 策略建議
    import re
    sections = _parse_md_sections(data.get('deep_analysis_md', ''))
    strategy_md = (sections.get('六、行銷策略建議', '')
                   or sections.get('六、行銷策略', ''))
    recs: List[str] = []
    if strategy_md:
        for m in re.finditer(r'(?:^|\n)\d+[.、]\s*\*{0,2}(.+?)\*{0,2}(?:\n|$)',
                             strategy_md):
            txt = m.group(1).strip()
            if txt and len(txt) > 5:
                recs.append(txt[:200])
    if not recs:
        recs = ['（請先完成 Stage 5 深度分析，AI 策略建議將自動填入）']

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # ── Slide 1：封面 ──────────────────────────────────
    _header_slide(
        prs,
        f'競品分析報告：{brand}',
        f'{purpose}  ｜  {industry} 產業  ｜  {service}',
        f'報告日期：{datetime.now().strftime("%Y年%m月%d日")}  ｜  目標受眾：{audience}'
    )

    # ── Slide 2：KPI 總覽 ──────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, BG)

    # 標題列
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.9), NAVY)
    _add_textbox(slide, Inches(0.4), Inches(0.1),
                 Inches(12), Inches(0.7),
                 '📊  調查數據總覽',
                 font_size=22, bold=True, color=WHITE)

    kpis = [
        ('競品數量', str(len(competitors))),
        ('分析文章', f'{total_cnt:,}'),
        ('追蹤關鍵字', str(kw_cnt)),
        ('監測平台', str(sum(1 for c in [dcard_cnt, ptt_cnt, news_cnt] if c > 0))),
    ]
    card_w  = Inches(2.8)
    card_h  = Inches(2.2)
    card_y  = Inches(1.5)
    gap     = Inches(0.5)
    start_x = Inches(0.8)

    for i, (label, val) in enumerate(kpis):
        cx = start_x + i * (card_w + gap)
        _add_rect(slide, cx, card_y, card_w, card_h, WHITE,
                  line_color=LIGHT_BLUE)
        # 頂部色條
        _add_rect(slide, cx, card_y, card_w, Inches(0.08), BLUE)
        # 數字
        _add_textbox(slide, cx, card_y + Inches(0.3),
                     card_w, Inches(1.1),
                     val, font_size=40, bold=True,
                     color=BLUE, align=PP_ALIGN.CENTER)
        # 標籤
        _add_textbox(slide, cx, card_y + Inches(1.5),
                     card_w, Inches(0.5),
                     label, font_size=13,
                     color=GRAY, align=PP_ALIGN.CENTER)

    # ── Slide 3：主要競品 ──────────────────────────────
    if competitors:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _solid_bg(slide, BG)
        _add_rect(slide, 0, 0, SLIDE_W, Inches(0.9), NAVY)
        _add_textbox(slide, Inches(0.4), Inches(0.1),
                     Inches(12), Inches(0.7),
                     f'🎯  主要競品（共 {len(competitors)} 個）',
                     font_size=22, bold=True, color=WHITE)

        # 每行3個，最多2行
        comp_w = Inches(3.8)
        comp_h = Inches(1.5)
        for idx, comp in enumerate(competitors[:6]):
            col = idx % 3
            row = idx // 3
            cx = Inches(0.4) + col * (comp_w + Inches(0.3))
            cy = Inches(1.1) + row * (comp_h + Inches(0.25))
            _add_rect(slide, cx, cy, comp_w, comp_h,
                      WHITE, line_color=LIGHT_BLUE)
            _add_rect(slide, cx, cy, comp_w, Inches(0.06), BLUE)

            name = comp.get('competitor_name', comp.get('name', f'競品{idx+1}'))
            label = (comp.get('label', comp.get('tags', '')) or '')[:25]
            kws   = ', '.join((comp.get('keywords', []) or [])[:3])

            _add_textbox(slide, cx + Inches(0.1), cy + Inches(0.1),
                         comp_w - Inches(0.2), Inches(0.5),
                         name, font_size=13, bold=True, color=NAVY)
            if label:
                _add_textbox(slide, cx + Inches(0.1), cy + Inches(0.55),
                             comp_w - Inches(0.2), Inches(0.35),
                             label, font_size=10, color=BLUE)
            if kws:
                _add_textbox(slide, cx + Inches(0.1), cy + Inches(0.95),
                             comp_w - Inches(0.2), Inches(0.45),
                             kws, font_size=9, color=GRAY)

    # ── Slide 4：社群聲量分布 ──────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, BG)
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.9), NAVY)
    _add_textbox(slide, Inches(0.4), Inches(0.1),
                 Inches(12), Inches(0.7),
                 '📊  社群聲量分布',
                 font_size=22, bold=True, color=WHITE)

    platform_data = [
        ('Dcard',      dcard_cnt, BLUE),
        ('PTT',        ptt_cnt,   RGBColor(0x27, 0xae, 0x60)),
        ('Google 新聞', news_cnt,  ORANGE),
    ]
    max_val = max(c for _, c, _ in platform_data) or 1
    bar_max_w = Inches(8)
    bar_h     = Inches(0.85)
    bar_y     = Inches(1.3)

    for i, (name, cnt, clr) in enumerate(platform_data):
        cy = bar_y + i * (bar_h + Inches(0.45))
        # 標籤
        _add_textbox(slide, Inches(0.5), cy,
                     Inches(1.8), bar_h,
                     name, font_size=14, bold=True, color=DARK)
        # 背景 track
        _add_rect(slide, Inches(2.5), cy + Inches(0.15),
                  bar_max_w, bar_h - Inches(0.3),
                  RGBColor(0xe8, 0xee, 0xf5))
        # 填色 bar
        bar_w = int(bar_max_w * cnt / max_val) if cnt else Inches(0.05)
        _add_rect(slide, Inches(2.5), cy + Inches(0.15),
                  bar_w, bar_h - Inches(0.3), clr)
        # 數字
        _add_textbox(slide, Inches(10.9), cy,
                     Inches(1.8), bar_h,
                     f'{cnt:,} 篇', font_size=14, bold=True, color=clr)

    # ── Slide 5：高互動文章 TOP 5 ──────────────────────
    if top_arts:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _solid_bg(slide, BG)
        _add_rect(slide, 0, 0, SLIDE_W, Inches(0.9), NAVY)
        _add_textbox(slide, Inches(0.4), Inches(0.1),
                     Inches(12), Inches(0.7),
                     '📰  高互動社群文章 TOP 5',
                     font_size=22, bold=True, color=WHITE)

        row_h = Inches(0.95)
        for idx, art in enumerate(top_arts[:5]):
            ry = Inches(1.0) + idx * row_h
            bg_clr = WHITE if idx % 2 == 0 else RGBColor(0xf8, 0xfb, 0xff)
            _add_rect(slide, Inches(0.3), ry,
                      SLIDE_W - Inches(0.6), row_h - Inches(0.05), bg_clr)

            eng = int(art['engagement'])
            eng_str = '爆' if eng >= 100 else str(eng)
            plat_clr = BLUE if art['platform'] == 'Dcard' else GREEN

            _add_textbox(slide, Inches(0.4), ry + Inches(0.15),
                         Inches(1.0), Inches(0.6),
                         art['platform'], font_size=10,
                         color=plat_clr, bold=True)
            _add_textbox(slide, Inches(1.5), ry + Inches(0.15),
                         Inches(9.5), Inches(0.65),
                         art['title'], font_size=12, color=DARK)
            _add_textbox(slide, Inches(11.2), ry + Inches(0.15),
                         Inches(1.7), Inches(0.65),
                         f'互動 {eng_str}', font_size=12,
                         bold=True, color=RED, align=PP_ALIGN.RIGHT)

    # ── Slide 6+：策略建議（每條一張） ────────────────
    _section_title_slide(prs, '💡', '行銷策略建議', BLUE)

    for i, rec in enumerate(recs[:5], 1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _solid_bg(slide, BG)

        # 編號圓圈背景
        _add_rect(slide, Inches(0.4), Inches(1.5),
                  Inches(1.2), Inches(1.2), BLUE)
        _add_textbox(slide, Inches(0.4), Inches(1.5),
                     Inches(1.2), Inches(1.2),
                     str(i), font_size=36, bold=True,
                     color=WHITE, align=PP_ALIGN.CENTER)

        # 建議內文
        _add_textbox(slide, Inches(2.0), Inches(1.2),
                     Inches(10.8), Inches(4.5),
                     rec, font_size=18, color=DARK, wrap=True)

        _add_textbox(slide, Inches(0.4), Inches(6.8),
                     Inches(12), Inches(0.4),
                     f'策略建議 {i} / {len(recs[:5])}  ｜  {brand} 競品分析',
                     font_size=10, color=GRAY)

    # ── 儲存 ──────────────────────────────────────────
    out = campaign_dir / 'report_executive_summary.pptx'
    prs.save(str(out))
    _safe_print(f'  ✅ 執行摘要 PPTX：{out.name}')
    return str(out)


# ═══════════════════════════════════════════════════════
#  PPTX 生成：競品對決版（Battlecard）
# ═══════════════════════════════════════════════════════

def generate_battlecard_pptx(campaign_dir: Path, data: Dict) -> str:
    """生成競品 Battlecard 投影片（可匯入 Canva / Google 簡報）。"""
    import re

    cfg = data.get('config', {})
    brand = cfg.get('brand', '我方品牌')
    service = cfg.get('service', '產品')
    competitors = data.get('competitors', [])

    sections = _parse_md_sections(data.get('deep_analysis_md', ''))
    battlecard_md = (sections.get('二、競品 Battlecard（前 3 名）', '')
                     or sections.get('二、競品 Battlecard', ''))

    def _extract_items(text: str) -> List[str]:
        items = []
        for line in text.split('\n'):
            m = re.match(r'^\s*[-*•]\s+(.+)', line)
            if m:
                items.append(m.group(1).strip()[:120])
        return items

    def _extract_sub(text: str, *keys: str) -> str:
        for key in keys:
            m = re.search(
                rf'(?:^|\n)\*{{0,2}}{re.escape(key)}\*{{0,2}}[：:]\s*(.*?)'
                rf'(?=\n\*{{0,2}}[^*\n]{{2,}}\*{{0,2}}[：:]|\Z)',
                text, re.DOTALL
            )
            if m:
                return m.group(1).strip()[:300]
        return ''

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # ── 封面 ──────────────────────────────────────────
    _header_slide(
        prs,
        f'競品對決 Battlecard：{brand}',
        f'核心產品：{service}  ｜  共分析 {len(competitors[:5])} 個競品',
        f'報告日期：{datetime.now().strftime("%Y年%m月%d日")}  ｜  競品對決版'
    )

    # ── 使用說明 ──────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, RGBColor(0xff, 0xfb, 0xf0))
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.9),
              ORANGE)
    _add_textbox(slide, Inches(0.4), Inches(0.1),
                 Inches(12), Inches(0.7),
                 '⚡  Battlecard 使用說明',
                 font_size=22, bold=True, color=WHITE)
    tips = [
        '🟢 我們贏的場景 → 在客戶面前主動提出，強化我方優勢',
        '🔴 我們輸的場景 → 預先準備話術，轉移焦點到我方優勢',
        '💬 建議話術 → 面對客戶詢問競品時的標準回應',
        '每競品一張投影片，可直接在 Canva / Google 簡報內自由編輯顏色與版面',
    ]
    for i, tip in enumerate(tips):
        _add_textbox(slide, Inches(0.8), Inches(1.2) + i * Inches(1.2),
                     Inches(12), Inches(1.0),
                     tip, font_size=15, color=DARK)

    # ── 每個競品一張 Battlecard ────────────────────────
    for idx, comp in enumerate(competitors[:5], 1):
        name = comp.get('competitor_name', comp.get('name', f'競品{idx}'))
        label = (comp.get('label', comp.get('tags', '')) or '')[:35]
        fan   = (comp.get('fan_profile', '') or '')[:80]
        kws   = comp.get('keywords', []) or []

        comp_md = ''
        if battlecard_md and name:
            import re as _re
            pat = rf'###\s*\d*\s*[.、]?\s*{_re.escape(name)}(.*?)(?=###\s*\d|\Z)'
            m = re.search(pat, battlecard_md, re.DOTALL | re.IGNORECASE)
            if m:
                comp_md = m.group(0)

        strengths  = _extract_items(_extract_sub(comp_md, '優勢', '強項'))[:5]
        weaknesses = _extract_items(_extract_sub(comp_md, '劣勢', '弱項'))[:5]
        our_adv    = _extract_items(_extract_sub(comp_md, '我方優勢', '我們的優勢'))[:5]
        win_sc     = _extract_sub(comp_md, '我們贏的場景', '贏的場景')[:180]
        lose_sc    = _extract_sub(comp_md, '我們輸的場景', '輸的場景')[:180]
        talking    = _extract_sub(comp_md, '話術建議', '建議話術', '異議回應')[:200]

        if not strengths and kws:
            strengths = [f'關鍵字佈局：{", ".join(kws[:3])}']
        if not strengths:
            strengths = ['（請完成 Stage 5 深度分析後自動填入）']
        if not our_adv:
            our_adv = ['（請完成 Stage 5 深度分析後自動填入）']

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _solid_bg(slide, BG)

        # 標題列（深藍）
        _add_rect(slide, 0, 0, SLIDE_W, Inches(0.95), NAVY)
        _add_textbox(slide, Inches(0.3), Inches(0.1),
                     Inches(8), Inches(0.75),
                     f'#{idx}  {name}',
                     font_size=22, bold=True, color=WHITE)
        if label:
            _add_textbox(slide, Inches(8.5), Inches(0.2),
                         Inches(4.5), Inches(0.55),
                         label, font_size=12, color=RGBColor(0xb8, 0xcf, 0xe8),
                         align=PP_ALIGN.RIGHT)

        # 左欄：競品分析（粉底）
        _add_rect(slide, 0, Inches(0.95),
                  Inches(6.3), SLIDE_H - Inches(0.95),
                  RGBColor(0xff, 0xf8, 0xf8))
        _add_textbox(slide, Inches(0.2), Inches(1.05),
                     Inches(5.9), Inches(0.4),
                     f'{name}  強項', font_size=11,
                     bold=True, color=GRAY)
        for i, s in enumerate(strengths):
            _add_textbox(slide, Inches(0.3), Inches(1.5) + i * Inches(0.5),
                         Inches(5.8), Inches(0.45),
                         f'✓  {s}', font_size=11, color=DARK)

        if weaknesses:
            sep_y = Inches(1.5) + len(strengths) * Inches(0.5) + Inches(0.1)
            _add_textbox(slide, Inches(0.2), sep_y,
                         Inches(5.9), Inches(0.4),
                         f'{name}  弱項', font_size=11,
                         bold=True, color=GRAY)
            for i, w in enumerate(weaknesses):
                _add_textbox(slide,
                             Inches(0.3), sep_y + Inches(0.45) + i * Inches(0.5),
                             Inches(5.8), Inches(0.45),
                             f'✗  {w}', font_size=11, color=RED)

        # 右欄：我方優勢（綠底）
        _add_rect(slide, Inches(6.5), Inches(0.95),
                  SLIDE_W - Inches(6.5), SLIDE_H - Inches(0.95),
                  RGBColor(0xf6, 0xff, 0xf8))
        _add_textbox(slide, Inches(6.6), Inches(1.05),
                     Inches(6.5), Inches(0.4),
                     f'{brand}  相對優勢', font_size=11,
                     bold=True, color=GRAY)
        for i, a in enumerate(our_adv):
            _add_textbox(slide, Inches(6.6), Inches(1.5) + i * Inches(0.5),
                         Inches(6.5), Inches(0.45),
                         f'✓  {a}', font_size=11, color=GREEN)

        if kws:
            _add_textbox(slide, Inches(6.6), Inches(4.2),
                         Inches(6.5), Inches(0.35),
                         '關鍵字', font_size=10, bold=True, color=GRAY)
            _add_textbox(slide, Inches(6.6), Inches(4.55),
                         Inches(6.5), Inches(0.6),
                         '  '.join(kws[:5]), font_size=10, color=BLUE)

        # 贏/輸場景（底部）
        if win_sc or lose_sc:
            _add_rect(slide, 0, Inches(5.45),
                      Inches(6.3), Inches(1.8),
                      RGBColor(0xe8, 0xf8, 0xf0))
            _add_textbox(slide, Inches(0.15), Inches(5.5),
                         Inches(6.0), Inches(0.4),
                         '🟢 我們贏的場景', font_size=11, bold=True, color=GREEN)
            _add_textbox(slide, Inches(0.15), Inches(5.9),
                         Inches(6.0), Inches(1.2),
                         win_sc or '（請完成深度分析後填入）',
                         font_size=10, color=DARK)

            _add_rect(slide, Inches(6.5), Inches(5.45),
                      SLIDE_W - Inches(6.5), Inches(1.8),
                      RGBColor(0xfe, 0xf0, 0xee))
            _add_textbox(slide, Inches(6.6), Inches(5.5),
                         Inches(6.5), Inches(0.4),
                         '🔴 我們輸的場景', font_size=11, bold=True, color=RED)
            _add_textbox(slide, Inches(6.6), Inches(5.9),
                         Inches(6.5), Inches(1.2),
                         lose_sc or '（請完成深度分析後填入）',
                         font_size=10, color=DARK)

        # 話術（底部）
        if talking:
            slide2 = prs.slides.add_slide(prs.slide_layouts[6])
            _solid_bg(slide2, RGBColor(0xf8, 0xfb, 0xff))
            _add_rect(slide2, 0, 0, SLIDE_W, Inches(0.9), BLUE)
            _add_textbox(slide2, Inches(0.3), Inches(0.1),
                         Inches(12), Inches(0.7),
                         f'💬  #{idx} {name} — 建議話術',
                         font_size=20, bold=True, color=WHITE)
            _add_textbox(slide2, Inches(0.8), Inches(1.2),
                         Inches(11.7), Inches(5.5),
                         talking, font_size=15, color=DARK, wrap=True)

    # ── 儲存 ──────────────────────────────────────────
    out = campaign_dir / 'report_battlecard.pptx'
    prs.save(str(out))
    _safe_print(f'  ✅ 競品對決 PPTX：{out.name}')
    return str(out)


# ═══════════════════════════════════════════════════════
#  Playwright PDF 生成
# ═══════════════════════════════════════════════════════

def generate_pdfs_with_playwright(html_files: Dict[str, str]) -> Dict[str, str]:
    """使用 Playwright 自動從 HTML 生成 PDF。"""
    try:
        from playwright.sync_api import sync_playwright
    except ImportError:
        print('  ⚠️ playwright 未安裝，跳過 PDF 生成')
        return {}

    labels = {
        'executive':    '執行摘要版',
        'battlecard':   '競品對決版',
        'deep_analysis': '深度分析版',
    }
    pdfs: Dict[str, str] = {}
    print('\n  🖨️  Playwright 生成 PDF...')

    try:
        with sync_playwright() as p:
            browser = p.chromium.launch(headless=True)
            for key, html_path in html_files.items():
                pdf_path = html_path.replace('.html', '.pdf')
                try:
                    page = browser.new_page()
                    uri = Path(html_path).absolute().as_uri()
                    page.goto(uri, wait_until='load')
                    page.wait_for_timeout(1000)
                    page.pdf(
                        path=pdf_path,
                        format='A4',
                        print_background=True,
                        margin={'top': '12mm', 'bottom': '12mm',
                                'left': '14mm', 'right': '14mm'},
                    )
                    page.close()
                    pdfs[key] = pdf_path
                    print(f'  ✅ {labels.get(key, key)} PDF：{Path(pdf_path).name}')
                except Exception as e:
                    print(f'  ⚠️ {labels.get(key, key)} PDF 失敗：{e}')
            browser.close()
    except Exception as e:
        print(f'  ⚠️ PDF 生成錯誤：{e}')

    return pdfs


# ═══════════════════════════════════════════════════════
#  主入口
# ═══════════════════════════════════════════════════════

def generate_pptx_reports(campaign_dir: str) -> Dict[str, str]:
    """生成 PPTX 報告（執行摘要 + Battlecard）。"""
    d = Path(campaign_dir)
    data = _load_data(d)
    outputs: Dict[str, str] = {}
    outputs['executive_pptx']  = generate_executive_pptx(d, data)
    outputs['battlecard_pptx'] = generate_battlecard_pptx(d, data)
    return outputs


if __name__ == '__main__':
    _fix_win_utf8()  # 只在直接執行時才包裝 stdout/stderr

    if len(sys.argv) < 2:
        print('用法：python pptx_report_generator.py <campaign_dir> [--pdf]')
        print('      --pdf 同時用 Playwright 生成 HTML → PDF')
        sys.exit(1)

    campaign_dir = sys.argv[1]
    want_pdf     = '--pdf' in sys.argv

    print('\n📊 PPTX 報告生成器啟動...')
    results = generate_pptx_reports(campaign_dir)

    if want_pdf:
        # 找同目錄的 HTML 報告
        html_files = {
            key.replace('_pptx', ''): path.replace('.pptx', '.html')
            for key, path in results.items()
            if Path(path.replace('.pptx', '.html')).exists()
        }
        if html_files:
            pdf_results = generate_pdfs_with_playwright(html_files)
            results.update({k + '_pdf': v for k, v in pdf_results.items()})

    print('\n✅ 產出檔案：')
    for key, path in results.items():
        print(f'  {key}: {Path(path).name}')

    print('\n📌 匯入說明：')
    print('  Canva   → 首頁 → 建立設計 → 匯入 → 選擇 .pptx → 可完整編輯')
    print('  Google  → slides.google.com → 檔案 → 匯入投影片 → 上傳 .pptx → 可完整編輯')
